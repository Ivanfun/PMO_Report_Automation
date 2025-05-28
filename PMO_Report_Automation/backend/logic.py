# PMO_Report_Automation/backend/logic.py

from docx import Document
from pptx import Presentation
from docx.shared import Cm, Pt, RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.text import WD_BREAK
import re
import os

# --- 輔助函式 ---

# 清理非法字元
def clean_text(text):
    if not isinstance(text, str):
        text = str(text)
    try:
        cleaned_text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\u200b\ufeff\u200c\u200d\u00ad]', '', text)
        return cleaned_text
    except Exception as e:
        print(f"清理文字時發生錯誤: {e} - 原始文字 (前100字): {text[:100]}")
        return ""

# 設定 Word 文件中文字區塊 (run) 的字型
def set_run_fonts(run, ascii_font, east_asia_font, hansi_font):
    try:
        rpr = run._element.get_or_add_rPr()
        rFonts_element = rpr.find(qn('w:rFonts'))
        if rFonts_element is not None:
            rpr.remove(rFonts_element)
        r_fonts = OxmlElement('w:rFonts')
        r_fonts.set(qn('w:ascii'), ascii_font)
        r_fonts.set(qn('w:eastAsia'), east_asia_font)
        r_fonts.set(qn('w:hAnsi'), hansi_font)
        rpr.append(r_fonts)
        run.font.name = ascii_font
    except Exception as e:
        print(f"設定文字區塊字型時發生錯誤: '{run.text[:50]}...'. 錯誤: {e}")

# 設定 Word 文件中 'Normal' 樣式的預設字型
def set_doc_normal_font(doc):
    try:
        style = doc.styles['Normal']
        font = style.font
        rpr = font._element.get_or_add_rPr()
        rFonts_element = rpr.find(qn('w:rFonts'))
        if rFonts_element is None:
            rFonts_element = OxmlElement('w:rFonts')
            rpr.append(rFonts_element)

        rFonts_element.set(qn('w:ascii'), 'Times New Roman')
        rFonts_element.set(qn('w:eastAsia'), '標楷體')
        rFonts_element.set(qn('w:hAnsi'), 'Times New Roman')
        font.name = 'Times New Roman'
        print("已成功設定 'Normal' 樣式的預設字型。")
    except Exception as e:
        print(f"警告: 無法設定 'Normal' 樣式的預設字型。錯誤: {e}")
        print("腳本將繼續執行，但預設樣式字型可能未正確應用。個別文字區塊的字型仍會設定。")

# 判斷一個表格行是否為空行
def is_row_empty(row_data):
    return all(not cell.strip() for cell in row_data)

# --- PPT 內容提取函式 ---

# 從 PPT 中提取文字列表內容
def extract_text_list_from_ppt(ppt, start_keyword, stop_keywords=None):
    items = []
    is_capturing = False
    if stop_keywords is None:
        stop_keywords = []

    for slide in ppt.slides:
        is_capturing = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            lines = re.split(r'[\n\r]+', text)
            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if start_keyword in line and not is_capturing:
                    is_capturing = True
                    continue

                if is_capturing:
                    if any(kw in line for kw in stop_keywords):
                        is_capturing = False
                        break
                    if line and line not in items:
                        items.append(line)
    return items

# 提取特定表格並將其內容結構化，同時處理專案階段表格的合併
def extract_relevant_tables_from_ppt(ppt):
    extracted_tables_data = {
        "work_summary_table_data": None,
        "project_stage_tables_data": []
    }

    for slide_idx, slide in enumerate(ppt.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_table:
                table = shape.table
                current_table_data = []
                for r in range(len(table.rows)):
                    row_texts = []
                    for c in range(len(table.columns)):
                        try:
                            cell_text = table.cell(r, c).text_frame.text.strip()
                            row_texts.append(clean_text(cell_text))
                        except Exception as e:
                            print(f"處理投影片 {slide_idx+1}、形狀 {shape_idx+1} 中的單元格 ({r},{c}) 時發生錯誤: {e}")
                            row_texts.append("")
                    current_table_data.append(row_texts)

                if not current_table_data:
                    continue

                first_cell_text = current_table_data[0][0].strip()

                if (("工作總覽" in first_cell_text or "時間" in first_cell_text) and
                        extracted_tables_data["work_summary_table_data"] is None):
                    extracted_tables_data["work_summary_table_data"] = current_table_data
                    print(f"    - 識別到 '工作總覽' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})")

                elif first_cell_text.startswith("專案階段"):
                    if extracted_tables_data["project_stage_tables_data"] and \
                       len(extracted_tables_data["project_stage_tables_data"][-1][0]) == len(current_table_data[0]):

                        if len(current_table_data) >= 2:
                             extracted_tables_data["project_stage_tables_data"][-1].extend(current_table_data[2:])
                             print(f"    - 合併 '專案階段' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})，並移除前兩行，到上一個表格")
                        else:
                             extracted_tables_data["project_stage_tables_data"][-1].extend(current_table_data[1:])
                             print(f"    - 合併 '專案階段' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})，並移除第一行，到上一個表格")
                    else:
                        extracted_tables_data["project_stage_tables_data"].append(current_table_data)
                        print(f"    - 識別到新的 '專案階段' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})")

    return extracted_tables_data

# 處理工作總覽表格數據，將其轉換為列表形式的關鍵字-值對
def process_work_summary_table(table_data):
    items = []
    if table_data:
        full_table_content = []
        for row_texts in table_data:
            if row_texts:
                full_table_content.append(" ".join(row_texts))
        combined_text = "\n".join(full_table_content)

        label_keywords = [
            "時間:", "階段:", "總體狀態:", "總體狀況:",
            "問題:", "風險:", "說明:",
            "本期重點:", "完成事項:", "待處理事項:", "下週計畫:", "重要進展:"
        ]
        split_pattern = r'(' + '|'.join(re.escape(kw) for kw in label_keywords) + r')'
        parts = re.split(split_pattern, combined_text)

        current_label = ""
        current_value = ""
        for part in parts:
            part = part.strip()
            if not part:
                continue
            if part in label_keywords:
                if current_label:
                    items.append(f"{current_label}{current_value.strip()}")
                current_label = part
                current_value = ""
            else:
                current_value += part + " "
        if current_label:
            items.append(f"{current_label}{current_value.strip()}")
        elif combined_text.strip() and not items:
            items.append(combined_text.strip())
    final_items = [item.strip() for item in items if item.strip()]
    return list(dict.fromkeys(final_items))

# --- Word 內容生成函式 ---

# 插入處理後的表格數據進 Word 並設定字型與欄寬、合併欄位與網底與換行支援
def add_filtered_tables(doc, tables_data_list):
    column_widths_cm = [1.24, 0.81, 0.81, 1.38, 4.69, 2.4, 2.4, 0.75, 2.27, 2.27]

    for table_data in tables_data_list:
        if not table_data:
            continue

        original_rows_count = len(table_data)
        last_valid_row_idx = -1
        for r_idx in range(len(table_data) - 1, -1, -1):
            if not is_row_empty(table_data[r_idx]):
                last_valid_row_idx = r_idx
                break

        if last_valid_row_idx == -1:
            table_data_to_use = []
        elif last_valid_row_idx < original_rows_count - 1:
            content_last_valid_row_idx = -1
            for r_idx in range(original_rows_count - 1, 0, -1):
                if not is_row_empty(table_data[r_idx]):
                    content_last_valid_row_idx = r_idx
                    break

            if original_rows_count > 0:
                if content_last_valid_row_idx == -1:
                    table_data_to_use = [table_data[0]]
                else:
                    table_data_to_use = table_data[:content_last_valid_row_idx + 1]
            else:
                table_data_to_use = []
        else:
            table_data_to_use = table_data

        if not table_data_to_use:
            print("    - 檢測到表格為空或僅包含空白行，已跳過。")
            continue

        rows = len(table_data_to_use)
        cols = len(table_data_to_use[0]) if rows > 0 else 0
        if cols == 0:
            print("    - 檢測到表格列數為零，已跳過。")
            continue

        doc_table = doc.add_table(rows=rows, cols=cols)
        doc_table.style = "Table Grid"

        for r in range(rows):
            for c in range(cols):
                cell_text = table_data_to_use[r][c] if c < len(table_data_to_use[r]) else ""

                word_cell = doc_table.cell(r, c)
                paragraph = word_cell.paragraphs[0]
                paragraph.clear()

                cleaned_text = clean_text(cell_text)

                if c == 4:
                    first_match = re.search(r'(\d+、)', cleaned_text)
                    if first_match:
                        temp_placeholder = "__FIRST_MATCH_PLACEHOLDER__"
                        temp_text = cleaned_text.replace(first_match.group(0), temp_placeholder, 1)
                        formatted_text = re.sub(r'(\d+、)', r'\n\1', temp_text)
                        formatted_text = formatted_text.replace(temp_placeholder, first_match.group(0))
                    else:
                        formatted_text = cleaned_text
                    lines_to_add = formatted_text.splitlines()

                elif c == 8 or c == 9:
                    original_lines = cleaned_text.splitlines()
                    lines_to_add = []
                    for idx, line in enumerate(original_lines):
                        if line.strip():
                            lines_to_add.append(line)
                            if idx < len(original_lines) - 1:
                                lines_to_add.append("")
                else:
                    lines_to_add = cleaned_text.splitlines()

                for idx, line in enumerate(lines_to_add):
                    if not line.strip() and (c == 8 or c == 9):
                        paragraph.add_run().add_break()
                        paragraph.add_run().add_break()
                    elif line.strip():
                        if idx > 0 and (c != 8 and c != 9):
                            paragraph.add_run().add_break()
                        elif idx > 0 and (c == 8 or c == 9) and lines_to_add[idx-1].strip() != "":
                            paragraph.add_run().add_break()

                        run = paragraph.add_run(line)
                        set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')
                        run.font.size = Pt(12)

                if r == 0:
                    tcPr = word_cell._tc.get_or_add_tcPr()
                    shd = OxmlElement('w:shd')
                    shd.set(qn('w:val'), 'clear')
                    shd.set(qn('w:color'), 'auto')
                    shd.set(qn('w:fill'), 'D9D9D9')
                    tcPr.append(shd)

        for col_idx in [0, 8, 9]:
            if col_idx < cols and rows > 1:
                try:
                    start_cell = doc_table.cell(1, col_idx)
                    end_cell = doc_table.cell(rows - 1, col_idx)
                    start_cell.merge(end_cell)
                except Exception as e:
                    print(f"無法合併表格第 {col_idx+1} 欄: {e}")

        for i, width_cm in enumerate(column_widths_cm):
            if i < cols:
                for cell in doc_table.columns[i].cells:
                    cell.width = Cm(width_cm)

# 插入 PPT 內容進 Word 文件中的動態會議區塊
def insert_dynamic_meeting_section(doc, ppt_name, work_items, summary_items, add_page_break=False):
    if add_page_break:
        doc.add_page_break()

    title = os.path.splitext(ppt_name)[0]

    heading2 = doc.add_heading(level=2)
    run_h2 = heading2.add_run(title)
    set_run_fonts(run_h2, 'Times New Roman', '標楷體', 'Times New Roman')

    heading3_work = doc.add_heading(level=3)
    run_h3_work = heading3_work.add_run("工作總覽")
    set_run_fonts(run_h3_work, 'Times New Roman', '標楷體', 'Times New Roman')

    if work_items:
        for item_text in work_items:
            cleaned_item_text = clean_text(item_text).replace('\r', '').replace('\n', ' ')
            para = doc.add_paragraph(f"● {cleaned_item_text}")
            para.paragraph_format.left_indent = Cm(0.63)
            for run in para.runs:
                set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')
    else:
        para_no_data = doc.add_paragraph("（尚無資料）")
        for run in para_no_data.runs:
            set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')

    heading3_summary = doc.add_heading(level=3)
    run_h3_summary = heading3_summary.add_run("本月概要")
    set_run_fonts(run_h3_summary, 'Times New Roman', '標楷體', 'Times New Roman')

    if summary_items:
        for item_text in summary_items:
            cleaned_item_text = clean_text(item_text).replace('\r', '').replace('\n', ' ')
            para = doc.add_paragraph(f"● {cleaned_item_text}")
            para.paragraph_format.left_indent = Cm(0.63)
            for run in para.runs:
                set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')
    else:
        para_no_data = doc.add_paragraph("（尚無資料）")
        for run in para_no_data.runs:
            set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')

# 新增函數：添加圖例和狀態說明
def add_legend_and_status(doc):
    p_legend = doc.add_paragraph()
    p_legend.paragraph_format.space_after = Pt(0)

    run_unit_title = p_legend.add_run("工項主要單位：")
    set_run_fonts(run_unit_title, 'Times New Roman', '標楷體', 'Times New Roman')
    run_unit_title.font.size = Pt(12)
    run_unit_title.font.color.rgb = RGBColor(0, 0, 0)

    run_ey_circle = p_legend.add_run(" ●")
    set_run_fonts(run_ey_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_ey_circle.font.size = Pt(12)
    run_ey_circle.font.color.rgb = RGBColor(0xF0, 0xD4, 0xB0)
    run_ey_text = p_legend.add_run(" EY")
    set_run_fonts(run_ey_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_ey_text.font.size = Pt(12)
    run_ey_text.font.color.rgb = RGBColor(0, 0, 0)

    run_vendor_circle = p_legend.add_run("  ●")
    set_run_fonts(run_vendor_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_vendor_circle.font.size = Pt(12)
    run_vendor_circle.font.color.rgb = RGBColor(0xF2, 0xAA, 0x6B)
    run_vendor_text = p_legend.add_run(" 建置廠商")
    set_run_fonts(run_vendor_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_vendor_text.font.size = Pt(12)
    run_vendor_text.font.color.rgb = RGBColor(0, 0, 0)

    run_first_ins_circle = p_legend.add_run("  ●")
    set_run_fonts(run_first_ins_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_first_ins_circle.font.size = Pt(12)
    run_first_ins_circle.font.color.rgb = RGBColor(0xC0, 0xEA, 0xFA)
    run_first_ins_text = p_legend.add_run(" 第一保")
    set_run_fonts(run_first_ins_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_first_ins_text.font.size = Pt(12)
    run_first_ins_text.font.color.rgb = RGBColor(0, 0, 0)

    run_status_title = p_legend.add_run("    狀態：")
    set_run_fonts(run_status_title, 'Times New Roman', '標楷體', 'Times New Roman')
    run_status_title.font.size = Pt(12)
    run_status_title.font.color.rgb = RGBColor(0, 0, 0)

    run_green_circle = p_legend.add_run(" ●")
    set_run_fonts(run_green_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_green_circle.font.size = Pt(12)
    run_green_circle.font.color.rgb = RGBColor(0x00, 0xB0, 0x50)
    run_green_text = p_legend.add_run(" 正常完成")
    set_run_fonts(run_green_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_green_text.font.size = Pt(12)
    run_green_text.font.color.rgb = RGBColor(0, 0, 0)

    run_yellow_circle = p_legend.add_run("  ●")
    set_run_fonts(run_yellow_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_yellow_circle.font.size = Pt(12)
    run_yellow_circle.font.color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    run_yellow_text = p_legend.add_run(" 部分延遲")
    set_run_fonts(run_yellow_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_yellow_text.font.size = Pt(12)
    run_yellow_text.font.color.rgb = RGBColor(0, 0, 0)

    run_red_circle = p_legend.add_run("  ●")
    set_run_fonts(run_red_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_red_circle.font.size = Pt(12)
    run_red_circle.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    run_red_text = p_legend.add_run(" 嚴重延遲")
    set_run_fonts(run_red_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_red_text.font.size = Pt(12)
    run_red_text.font.color.rgb = RGBColor(0, 0, 0)