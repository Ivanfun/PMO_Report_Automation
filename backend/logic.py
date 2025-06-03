# PMO_Report_Automation/backend/logic.py

from docx import Document
from pptx import Presentation
from docx.shared import Cm, Pt, RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.text import WD_BREAK
import re
import os

# --- 輔助函式 ---

# 清理非法字元，並增加將全形空格轉換為半形空格的處理
def clean_text(text):
    if not isinstance(text, str):
        text = str(text)
    try:
        # 移除控制字符和零寬度字符
        cleaned_text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\u200b\ufeff\u200c\u200d\u00ad]', '', text)
        # 將全形空格轉換為半形空格
        cleaned_text = cleaned_text.replace('\u3000', ' ')
        return cleaned_text
    except Exception as e:
        print(f"清理文字時發生錯誤: {e} - 原始文字 (前100字): {text[:100]}")
        return ""

# 設定 Word 文件中文字區塊 (run) 的字型
def set_run_fonts(run, ascii_font, east_asia_font, hansi_font):
    try:
        rpr = run._element.get_or_add_rPr()
        rFonts_element = rpr.find(qn('w:rFonts'))
        if rFonts_element is not None:
            rpr.remove(rFonts_element)
        r_fonts = OxmlElement('w:rFonts')
        r_fonts.set(qn('w:ascii'), ascii_font)
        r_fonts.set(qn('w:eastAsia'), east_asia_font)
        r_fonts.set(qn('w:hAnsi'), hansi_font)
        rpr.append(r_fonts)
        run.font.name = ascii_font # 再次設置 run.font.name 確保應用
    except Exception as e:
        print(f"設定文字區塊字型時發生錯誤: '{run.text[:50]}...'. 錯誤: {e}")

# 設定 Word 文件中 'Normal' 樣式的預設字型
def set_doc_normal_font(doc):
    try:
        style = doc.styles['Normal']
        font = style.font
        rpr = font._element.get_or_add_rPr()
        rFonts_element = rpr.find(qn('w:rFonts'))
        if rFonts_element is None:
            rFonts_element = OxmlElement('w:rFonts')
            rpr.append(rFonts_element)

        rFonts_element.set(qn('w:ascii'), 'Times New Roman')
        rFonts_element.set(qn('w:eastAsia'), '標楷體')
        rFonts_element.set(qn('w:hAnsi'), 'Times New Roman')
        font.name = 'Times New Roman' # 再次設置 font.name 確保應用
        print("已成功設定 'Normal' 樣式的預設字型。")
    except Exception as e:
        print(f"警告: 無法設定 'Normal' 樣式的預設字型。錯誤: {e}")
        print("腳本將繼續執行，但預設樣式字型可能未正確應用。個別文字區塊的字型仍會設定。")

# 判斷一個表格行是否為空行
def is_row_empty(row_data):
    return all(not cell.strip() for cell in row_data)

# --- PPT 內容提取函式 ---

# 從 PPT 中提取文字列表內容 (如「本月概要」)
def extract_text_list_from_ppt(ppt, start_keyword, stop_keywords=None):
    items = []
    is_capturing = False
    if stop_keywords is None:
        stop_keywords = []

    for slide in ppt.slides:
        is_capturing_on_this_slide = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            lines = re.split(r'[\n\r]+', text)
            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if start_keyword in line and not is_capturing_on_this_slide:
                    is_capturing_on_this_slide = True
                    is_capturing = True
                    continue

                if is_capturing_on_this_slide:
                    if any(kw in line for kw in stop_keywords):
                        is_capturing_on_this_slide = False
                        break
                    
                    if line and line not in items:
                        items.append(line)
        
        if not is_capturing_on_this_slide and is_capturing:
            break
    return items

# 提取特定表格並將其內容結構化，同時處理專案階段表格的合併
def extract_relevant_tables_from_ppt(ppt):
    extracted_tables_data = {
        "work_summary_table_data": None,
        "project_stage_tables_data": [],
        "program_modifications_table_data": [],
        "program_modifications_slide_indices": []
    }

    # 定義目標關鍵字列表
    target_keywords = [
        "修改因重跑過程發現的問題所產生的程式修改",
        "修改項目" # 根據您的需求增加「修改項目」作為關鍵字
    ]
    # 將關鍵字標準化，用於比對
    standardized_target_keywords = [clean_text(kw).replace(" ", "") for kw in target_keywords]

    for slide_idx, slide in enumerate(ppt.slides): # slide_idx 從 0 開始
        print(f"\n--- 正在檢查投影片 {slide_idx+1} ---")

        # 提取整個投影片的文本內容進行關鍵字檢查
        slide_full_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_full_text += clean_text(shape.text_frame.text) + "\n"
        
        # 標準化投影片的文本，移除空格進行更寬鬆的比對
        standardized_slide_full_text = slide_full_text.replace(" ", "")

        # 檢查關鍵字是否存在於整個投影片文本中
        is_keyword_present_in_slide_anywhere = False
        detected_keyword = None
        for kw in standardized_target_keywords:
            if kw in standardized_slide_full_text:
                is_keyword_present_in_slide_anywhere = True
                detected_keyword = kw # 記錄是哪個關鍵字被偵測到
                print(f"  ✅ 投影片 {slide_idx+1} 的整體文本中偵測到關鍵字: '{kw}'")
                break
        
        if not is_keyword_present_in_slide_anywhere:
            print(f"  ❌ 投影片 {slide_idx+1} 的整體文本中未偵測到任何關鍵字。")

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_table:
                table = shape.table
                current_table_data = []
                for r in range(len(table.rows)):
                    row_texts = []
                    for c in range(len(table.columns)):
                        try:
                            cell_text = table.cell(r, c).text_frame.text.strip()
                            row_texts.append(clean_text(cell_text))
                        except Exception as e:
                            print(f"    處理投影片 {slide_idx+1}、形狀 {shape_idx+1} 中的單元格 ({r},{c}) 時發生錯誤: {e}")
                            row_texts.append("")
                    current_table_data.append(row_texts)

                if not current_table_data:
                    print(f"    形狀 {shape_idx+1} 中的表格為空，跳過。")
                    continue

                first_cell_text = current_table_data[0][0].strip()
                # 標準化表格第一格文本
                standardized_first_cell_text = clean_text(first_cell_text).replace(" ", "")
                print(f"    形狀 {shape_idx+1} 中的表格第一格內容: '{first_cell_text}'")

                # 識別「工作總覽」表格
                if (("工作總覽" in standardized_first_cell_text or "時間" in standardized_first_cell_text) and
                        extracted_tables_data["work_summary_table_data"] is None):
                    extracted_tables_data["work_summary_table_data"] = current_table_data
                    print(f"    - 識別到 '工作總覽' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})")

                # 識別並合併「專案階段」表格
                elif standardized_first_cell_text.startswith("專案階段"):
                    if extracted_tables_data["project_stage_tables_data"] and \
                       len(extracted_tables_data["project_stage_tables_data"][-1][0]) == len(current_table_data[0]):
                        
                        if len(current_table_data) >= 2 and \
                           all(clean_text(current_table_data[1][c]).replace(" ", "") == clean_text(extracted_tables_data["project_stage_tables_data"][-1][1][c]).replace(" ", "") for c in range(len(current_table_data[1]))):
                            extracted_tables_data["project_stage_tables_data"][-1].extend(current_table_data[2:])
                            print(f"    - 合併 '專案階段' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})，並移除前兩行，到上一個表格")
                        else:
                            extracted_tables_data["project_stage_tables_data"][-1].extend(current_table_data[1:])
                            print(f"    - 合併 '專案階段' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})，並移除第一行，到上一個表格")
                    else:
                        extracted_tables_data["project_stage_tables_data"].append(current_table_data)
                        print(f"    - 識別到新的 '專案階段' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})")

                # 識別「程式修改」表格（條件放寬：只要整頁有關鍵字，不再檢查表格第一欄，也不限第二頁）
                elif is_keyword_present_in_slide_anywhere:
                    extracted_tables_data["program_modifications_table_data"].append(current_table_data)
                    extracted_tables_data["program_modifications_slide_indices"].append(slide_idx)
                    print(f"    - 偵測到符合條件的 '程式修改' 表格 (投影片 {slide_idx+1}, 形狀 {shape_idx+1})。")

                else:
                    print(f"    - 形狀 {shape_idx+1} 中的表格不是目標表格類型。")
    return extracted_tables_data

# 處理工作總覽表格數據，將其轉換為列表形式的關鍵字-值對
def process_work_summary_table(table_data):
    items = []
    if table_data:
        full_table_content = []
        for row_texts in table_data:
            if any(cell.strip() for cell in row_texts):
                full_table_content.append(" ".join(row_texts).strip())

        combined_text = "\n".join(full_table_content)

        label_keywords = [
            "時間:", "階段:", "總體狀態:", "總體狀況:",
            "問題:", "風險:", "說明:",
            "本期重點:", "完成事項:", "待處理事項:", "下週計畫:", "重要進展:"
        ]
        split_pattern = r'(' + '|'.join(re.escape(kw) for kw in label_keywords) + r')'
        parts = re.split(split_pattern, combined_text)

        current_label = ""
        current_value = ""
        for part in parts:
            part = part.strip()
            if not part:
                continue

            if part in label_keywords:
                if current_label:
                    items.append(f"{current_label}{current_value.strip()}")
                current_label = part
                current_value = ""
            else:
                current_value += part + " "

        if current_label:
            items.append(f"{current_label}{current_value.strip()}")
        elif combined_text.strip() and not items:
            items.append(combined_text.strip())

    final_items = [item.strip() for item in items if item.strip()]
    return list(dict.fromkeys(final_items))

# 插入處理後的表格數據進 Word 並設定字型與欄寬、合併欄位與網底與換行支援
def add_filtered_tables(doc, tables_data_list):
    predefined_widths_map = {
        10: [Cm(1.24), Cm(0.81), Cm(0.81), Cm(1.38), Cm(4.69), Cm(2.56), Cm(2.25), Cm(0.75), Cm(2.25), Cm(2.29)],
        2: [Cm(2), Cm(17.03)],
    }

    for table_data in tables_data_list:
        if not table_data:
            continue

        # 過濾空白行
        last_valid_row_idx = -1
        for r_idx in range(len(table_data) - 1, -1, -1):
            if not is_row_empty(table_data[r_idx]):
                last_valid_row_idx = r_idx
                break

        table_data_to_use = table_data[:last_valid_row_idx + 1] if last_valid_row_idx != -1 else []
        if not table_data_to_use:
            print("    - 檢測到表格為空或僅包含空白行，已跳過。")
            continue

        rows = len(table_data_to_use)
        cols = len(table_data_to_use[0]) if rows > 0 else 0
        if cols == 0:
            print("    - 檢測到表格列數為零，已跳過。")
            continue

        # 建立表格
        doc_table = doc.add_table(rows=rows, cols=cols)
        doc_table.style = "Table Grid"
        doc_table.autofit = False

        tbl = doc_table._tbl
        tblPr = tbl.tblPr
        if tblPr is None:
            tblPr = OxmlElement('w:tblPr')
            tbl.insert(0, tblPr)

        tblW_existing = tblPr.find(qn('w:tblW'))
        if tblW_existing is not None:
            tblPr.remove(tblW_existing)

        tblW = OxmlElement('w:tblW')
        tblW.set(qn('w:type'), 'dxa')
        tblW.set(qn('w:w'), str(int(19.03 * 567)))  # 11361 dxa
        tblPr.append(tblW)

        # 欄寬設定
        if cols in predefined_widths_map:
            current_column_widths_cm = predefined_widths_map[cols]
            print(f"    - 為 {cols} 欄表格套用預設欄寬: {current_column_widths_cm}")
        else:
            col_width = Cm(19.03 / cols)
            current_column_widths_cm = [col_width] * cols
            print(f"    - 為 {cols} 欄表格平均分配欄寬: {col_width.cm:.2f} cm/欄")

        # 寫入內容
        for r in range(rows):
            for c in range(cols):
                cell_text = table_data_to_use[r][c] if c < len(table_data_to_use[r]) else ""
                word_cell = doc_table.cell(r, c)
                paragraph = word_cell.paragraphs[0]
                paragraph.clear()
                cleaned_text = clean_text(cell_text)

                # 特殊換行處理（第4欄、第9~10欄）
                if c == 4:
                    first_match = re.search(r'(\d+、)', cleaned_text)
                    if first_match:
                        temp_placeholder = "__FIRST_MATCH_PLACEHOLDER__"
                        temp_text = cleaned_text.replace(first_match.group(0), temp_placeholder, 1)
                        formatted_text = re.sub(r'(\d+、)', r'\n\1', temp_text)
                        formatted_text = formatted_text.replace(temp_placeholder, first_match.group(0))
                    else:
                        formatted_text = cleaned_text
                    lines_to_add = formatted_text.splitlines()
                elif c in [8, 9]:
                    lines_to_add = cleaned_text.splitlines()

                else:
                    lines_to_add = cleaned_text.splitlines()

                for idx, line in enumerate(lines_to_add):
                    if not line.strip():
                        continue  # 跳過空行

                    elif line.strip():
                        if idx > 0 and (c not in [8, 9] or lines_to_add[idx-1].strip()):
                            paragraph.add_run().add_break()
                        run = paragraph.add_run(line)
                        set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')
                        run.font.size = Pt(12)

                # 表頭網底
                if r == 0:
                    tcPr = word_cell._tc.get_or_add_tcPr()
                    shd = OxmlElement('w:shd')
                    shd.set(qn('w:val'), 'clear')
                    shd.set(qn('w:color'), 'auto')
                    shd.set(qn('w:fill'), 'D9D9D9')  # 淺灰色
                    tcPr.append(shd)

        # 合併欄位（僅適用 10 欄的表格）
        if cols == 10 and rows > 1:
            for col_idx in [0, 8, 9]:
                if col_idx < cols:
                    try:
                        doc_table.cell(1, col_idx).merge(doc_table.cell(rows - 1, col_idx))
                    except Exception as e:
                        print(f"無法合併表格第 {col_idx+1} 欄: {e}")

        # 設定每一欄的寬度
        for i, width_cm in enumerate(current_column_widths_cm):
            if i < cols:
                for cell in doc_table.columns[i].cells:
                    cell.width = width_cm



# 插入 PPT 內容進 Word 文件中的動態會議區塊
def insert_dynamic_meeting_section(doc, ppt_name, work_items, summary_items, add_page_break=False):
    if add_page_break:
        doc.add_page_break()

    title = os.path.splitext(ppt_name)[0]

    heading2 = doc.add_heading(level=2)
    run_h2 = heading2.add_run(title)
    set_run_fonts(run_h2, 'Times New Roman', '標楷體', 'Times New Roman')

    heading3_work = doc.add_heading(level=3)
    run_h3_work = heading3_work.add_run("工作總覽")
    set_run_fonts(run_h3_work, 'Times New Roman', '標楷體', 'Times New Roman')

    if work_items:
        for item_text in work_items:
            cleaned_item_text = clean_text(item_text).replace('\r', '').replace('\n', ' ')
            para = doc.add_paragraph(f"● {cleaned_item_text}")
            para.paragraph_format.left_indent = Cm(0.63)
            for run in para.runs:
                set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')
    else:
        para_no_data = doc.add_paragraph("（尚無資料）")
        for run in para_no_data.runs:
            set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')

    heading3_summary = doc.add_heading(level=3)
    run_h3_summary = heading3_summary.add_run("本月概要")
    set_run_fonts(run_h3_summary, 'Times New Roman', '標楷體', 'Times New Roman')

    if summary_items:
        for item_text in summary_items:
            cleaned_item_text = clean_text(item_text).replace('\r', '').replace('\n', ' ')
            para = doc.add_paragraph(f"● {cleaned_item_text}")
            para.paragraph_format.left_indent = Cm(0.63)
            for run in para.runs:
                set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')
    else:
        para_no_data = doc.add_paragraph("（尚無資料）")
        for run in para_no_data.runs:
            set_run_fonts(run, 'Times New Roman', '標楷體', 'Times New Roman')

# 新增函數：添加圖例和狀態說明
def add_legend_and_status(doc):
    p_legend = doc.add_paragraph()
    p_legend.paragraph_format.space_after = Pt(0)

    run_unit_title = p_legend.add_run("工項主要單位：")
    set_run_fonts(run_unit_title, 'Times New Roman', '標楷體', 'Times New Roman')
    run_unit_title.font.size = Pt(12)
    run_unit_title.font.color.rgb = RGBColor(0, 0, 0)

    run_ey_circle = p_legend.add_run(" ●")
    set_run_fonts(run_ey_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_ey_circle.font.size = Pt(12)
    run_ey_circle.font.color.rgb = RGBColor(0xF0, 0xD4, 0xB0)
    run_ey_text = p_legend.add_run(" EY")
    set_run_fonts(run_ey_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_ey_text.font.size = Pt(12)
    run_ey_text.font.color.rgb = RGBColor(0, 0, 0)

    run_vendor_circle = p_legend.add_run("  ●")
    set_run_fonts(run_vendor_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_vendor_circle.font.size = Pt(12)
    run_vendor_circle.font.color.rgb = RGBColor(0xF2, 0xAA, 0x6B)
    run_vendor_text = p_legend.add_run(" 建置廠商")
    set_run_fonts(run_vendor_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_vendor_text.font.size = Pt(12)
    run_vendor_text.font.color.rgb = RGBColor(0, 0, 0)

    run_first_ins_circle = p_legend.add_run("  ●")
    set_run_fonts(run_first_ins_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_first_ins_circle.font.size = Pt(12)
    run_first_ins_circle.font.color.rgb = RGBColor(0xC0, 0xEA, 0xFA)
    run_first_ins_text = p_legend.add_run(" 第一保")
    set_run_fonts(run_first_ins_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_first_ins_text.font.size = Pt(12)
    run_first_ins_text.font.color.rgb = RGBColor(0, 0, 0)

    run_status_title = p_legend.add_run("    狀態：")
    set_run_fonts(run_status_title, 'Times New Roman', '標楷體', 'Times New Roman')
    run_status_title.font.size = Pt(12)
    run_status_title.font.color.rgb = RGBColor(0, 0, 0)

    run_green_circle = p_legend.add_run(" ●")
    set_run_fonts(run_green_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_green_circle.font.size = Pt(12)
    run_green_circle.font.color.rgb = RGBColor(0x00, 0xB0, 0x50)
    run_green_text = p_legend.add_run(" 正常完成")
    set_run_fonts(run_green_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_green_text.font.size = Pt(12)
    run_green_text.font.color.rgb = RGBColor(0, 0, 0)

    run_yellow_circle = p_legend.add_run("  ●")
    set_run_fonts(run_yellow_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_yellow_circle.font.size = Pt(12)
    run_yellow_circle.font.color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    run_yellow_text = p_legend.add_run(" 部分延遲")
    set_run_fonts(run_yellow_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_yellow_text.font.size = Pt(12)
    run_yellow_text.font.color.rgb = RGBColor(0, 0, 0)

    run_red_circle = p_legend.add_run("  ●")
    set_run_fonts(run_red_circle, 'Times New Roman', '標楷體', 'Times New Roman')
    run_red_circle.font.size = Pt(12)
    run_red_circle.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    run_red_text = p_legend.add_run(" 嚴重延遲")
    set_run_fonts(run_red_text, 'Times New Roman', '標楷體', 'Times New Roman')
    run_red_text.font.size = Pt(12)
    run_red_text.font.color.rgb = RGBColor(0, 0, 0)